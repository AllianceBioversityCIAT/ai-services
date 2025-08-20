import docx
import boto3
import pandas as pd
from io import BytesIO
from PyPDF2 import PdfReader
from pptx import Presentation
from app.utils.config.config_util import S3
from app.utils.logger.logger_util import get_logger

logger = get_logger()

s3_client = boto3.client(
    's3',
    aws_access_key_id=S3['aws_access_key'],
    aws_secret_access_key=S3['aws_secret_key'],
    region_name=S3['aws_region']
)


def _process_file_content(file_extension, file_content):
    if file_extension == 'pdf':
        logger.info("📄 Processing PDF file...")
        pdf_reader = PdfReader(BytesIO(file_content))
        text = "".join(page.extract_text() + "\n" for page in pdf_reader.pages)
        return text
    elif file_extension == 'docx':
        logger.info("📄 Processing DOCX file...")
        doc = docx.Document(BytesIO(file_content))
        text = "".join(para.text + "\n" for para in doc.paragraphs)
        return text
    elif file_extension == 'txt':
        logger.info("📄 Processing TXT file...")
        return file_content.decode('utf-8')
    elif file_extension in ('xls', 'xlsx'):
        logger.info("📄 Processing EXCEL file...")
        df = pd.read_excel(BytesIO(file_content), header=0)
        logger.info(f"📊 Original DataFrame shape: {df.shape}")
        
        df = df.dropna(axis=1, how='all')        
        df = df.dropna(axis=0, how='all')
        df = df[~df.apply(lambda row: all(str(val).strip() == '' or pd.isna(val) for val in row), axis=1)]
        df = df.drop_duplicates()
        df = df.reset_index(drop=True)
        logger.info(f"📊 Cleaned DataFrame shape: {df.shape}")
    
        try:
            structured_rows = []
            for index, row in df.iterrows():
                row_parts = []
                for col in df.columns:
                    value = str(row[col]).strip()
                    if value and value != 'nan' and value != 'None':
                        row_parts.append(f"{col}: {value}")
                
                if row_parts:
                    row_text = ", ".join(row_parts)
                    structured_rows.append(row_text)
            
            logger.info(f"📊 Processed {len(structured_rows)} meaningful Excel rows as individual chunks")
            
            if structured_rows:
                logger.info("📝 Sample rows:")
                for i, row in enumerate(structured_rows[:3]):
                    logger.info(f"  Row {i+1}: {row}")
            
            return {"type": "excel", "chunks": structured_rows}
            
        except Exception as e:
            logger.warning(f"⚠️ Excel processing failed, falling back to CSV: {e}")
            df = df.to_csv(index=False, header=True)
            return df

    elif file_extension == 'pptx':
        logger.info("📄 Processing PPTX file...")
        prs = Presentation(BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError(f"❌ File format not supported: {file_extension}")


def read_document_from_s3(bucket_name, file_key):
    try:
        logger.info(
            f"📂 Downloading the {file_key} file from the bucket {bucket_name}...")
        response = s3_client.get_object(Bucket=bucket_name, Key=file_key)
        file_content = response['Body'].read()
        file_extension = file_key.lower().split('.')[-1]

        return _process_file_content(file_extension, file_content)

    except Exception as e:
        logger.error(
            f"❌ Error while reading {file_key} from bucket {bucket_name}: {str(e)}")
        raise


def upload_file_to_s3(file_content, bucket_name, file_key, content_type=None):
    """
    Upload a file to an S3 bucket.

    Args:
        file_content: The content of the file to upload (bytes or file-like object)
        bucket_name: Name of the S3 bucket
        file_key: The key (path) where the file will be stored in S3
        content_type: Optional MIME type of the file

    Returns:
        dict: The response from S3 upload operation

    Raises:
        Exception: If the upload fails
    """
    try:
        logger.info(f"📤 Uploading file to {bucket_name}/{file_key}...")

        # Prepare upload parameters
        upload_args = {
            'Bucket': bucket_name,
            'Key': file_key,
            'Body': file_content
        }

        # Add content type if provided
        if content_type:
            upload_args['ContentType'] = content_type

        # Upload the file
        response = s3_client.put_object(**upload_args)

        logger.info(
            f"✅ File successfully uploaded to {bucket_name}/{file_key}")
        return response

    except Exception as e:
        logger.error(
            f"❌ Error uploading file to {bucket_name}/{file_key}: {str(e)}")
        raise
