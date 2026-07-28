import json
import boto3
import mammoth
import html2text
import pandas as pd
from io import BytesIO
from pptx import Presentation
from botocore.exceptions import ClientError
from utils.logger.logger_util import get_logger
from utils.config.config_util import get_boto3_client_kwargs

logger = get_logger()

_html_to_text = html2text.HTML2Text()
_html_to_text.body_width = 0
_html_to_text.ignore_links = True
_html_to_text.ignore_images = True

SUPPORTED_DOCUMENT_EXTENSIONS = {
    'pdf', 'jpg', 'jpeg', 'png', 'tiff', 'tif',
    'docx', 'txt', 'xls', 'xlsx', 'pptx',
}
MAX_PROJECT_DOCUMENTS = 3
RESPONSE_JSON_FILENAME = "response.json"

_s3_client = None


def _get_s3_client():
    """Lazy client — avoids caching invalid keys across warm Lambda containers."""
    global _s3_client
    if _s3_client is None:
        _s3_client = boto3.client("s3", **get_boto3_client_kwargs())
    return _s3_client


def _process_file_content(file_extension, file_content):
    if file_extension == 'docx':
        logger.info("📄 Processing DOCX file with Mammoth...")
        result = mammoth.convert_to_html(BytesIO(file_content))
        for message in result.messages:
            logger.warning(f"Mammoth conversion warning: {message}")
        html = result.value.strip()
        text = _html_to_text.handle(html).strip()
        logger.info(f"✅ DOCX converted to plain text — {len(text)} characters extracted")
        return text
    elif file_extension == 'txt':
        logger.info("📄 Processing TXT file...")
        return file_content.decode('utf-8')
    elif file_extension in ('xls', 'xlsx'):
        logger.info("📄 Processing EXCEL file...")
        df = pd.read_excel(BytesIO(file_content), header=0)
        logger.info(f"📊 Original DataFrame shape: {df.shape}")
        
        df = df.dropna(axis=1, how='all')        
        df = df.dropna(axis=0, how='all')
        df = df[~df.apply(lambda row: all(str(val).strip() == '' or pd.isna(val) for val in row), axis=1)]
        df = df.drop_duplicates()
        df = df.reset_index(drop=True)
        logger.info(f"📊 Cleaned DataFrame shape: {df.shape}")
    
        try:
            structured_rows = []
            for index, row in df.iterrows():
                row_parts = []
                for col in df.columns:
                    value = str(row[col]).strip()
                    if value and value != 'nan' and value != 'None':
                        row_parts.append(f"{col}: {value}")
                
                if row_parts:
                    row_text = ", ".join(row_parts)
                    structured_rows.append(row_text)
            
            logger.info(f"📊 Processed {len(structured_rows)} meaningful Excel rows as individual chunks")
            
            if structured_rows:
                logger.info("📝 Sample rows:")
                for i, row in enumerate(structured_rows[:3]):
                    logger.info(f"  Row {i+1}: {row}")
            
            return {"type": "excel", "chunks": structured_rows}
            
        except Exception as e:
            logger.warning(f"⚠️ Excel processing failed, falling back to CSV: {e}")
            df = df.to_csv(index=False, header=True)
            return df

    elif file_extension == 'pptx':
        logger.info("📄 Processing PPTX file...")
        prs = Presentation(BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError(f"File format not supported: {file_extension}")


def _normalize_project_folder(project_folder: str) -> str:
    return project_folder.strip('/')


def get_response_json_key(project_folder: str) -> str:
    """Build the S3 key for the cached project overview response."""
    prefix = _normalize_project_folder(project_folder)
    return f"{prefix}/{RESPONSE_JSON_FILENAME}" if prefix else RESPONSE_JSON_FILENAME


def list_available_project_files(
    bucket_name: str,
    project_folder: str,
) -> list[dict]:
    """
    List files currently available in a project folder (excluding response.json).

    Returns:
        Sorted list of dicts with file_key and file_name. Empty list if none found.
    """
    prefix = _normalize_project_folder(project_folder)
    if prefix:
        prefix = f"{prefix}/"

    logger.info(f"📂 Listing available files in s3://{bucket_name}/{prefix}")

    response = _get_s3_client().list_objects_v2(Bucket=bucket_name, Prefix=prefix)
    contents = response.get('Contents', [])

    files = []
    for obj in contents:
        key = obj['Key']
        if key.endswith('/') or obj.get('Size', 0) == 0:
            continue

        file_name = key.rsplit('/', 1)[-1]
        if file_name == RESPONSE_JSON_FILENAME:
            continue

        files.append({
            "file_key": key,
            "file_name": file_name,
        })

    files.sort(key=lambda item: item["file_name"])
    logger.info(f"📂 Found {len(files)} available file(s) in project folder")
    return files


def list_project_documents(
    bucket_name: str,
    project_folder: str,
    max_documents: int = MAX_PROJECT_DOCUMENTS,
) -> list[str]:
    """
    List supported document keys under a project folder in S3.

    Args:
        bucket_name: S3 bucket name
        project_folder: Folder prefix for the project (e.g. star/ai-insights/projects/abc123)
        max_documents: Maximum number of documents allowed (default: 3)

    Returns:
        Sorted list of S3 object keys

    Raises:
        ValueError: If no supported documents are found or the limit is exceeded
    """
    available_files = list_available_project_files(bucket_name, project_folder)

    document_keys = []
    for file_info in available_files:
        key = file_info["file_key"]
        extension = key.lower().rsplit('.', 1)[-1]
        if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
            logger.warning(f"Skipping unsupported file: {key}")
            continue
        document_keys.append(key)

    if not document_keys:
        raise ValueError(
            f"No supported documents found in s3://{bucket_name}/{project_folder}"
        )

    if len(document_keys) > max_documents:
        raise ValueError(
            f"Found {len(document_keys)} documents in project folder "
            f"(maximum allowed: {max_documents})"
        )

    logger.info(f"📄 Found {len(document_keys)} document(s): {document_keys}")
    return document_keys


def delete_project_files(
    bucket_name: str,
    project_folder: str,
    file_names: list[str],
) -> list[str]:
    """
    Delete one or more files from a project folder in S3.

    Only files inside the given project folder can be deleted.
    response.json cannot be deleted through this function.

    Args:
        bucket_name: S3 bucket name
        project_folder: Folder prefix for the project
        file_names: File names (not full keys) to delete

    Returns:
        List of deleted file names

    Raises:
        ValueError: If file_names is empty or includes invalid names
    """
    if not file_names:
        raise ValueError("At least one file_name is required for deletion")

    prefix = _normalize_project_folder(project_folder)
    deleted = []

    for file_name in file_names:
        clean_name = file_name.strip().lstrip('/')
        if not clean_name or '/' in clean_name:
            raise ValueError(
                f"Invalid file_name '{file_name}'. Provide only the file name, not a path."
            )
        if clean_name == RESPONSE_JSON_FILENAME:
            raise ValueError(f"Cannot delete '{RESPONSE_JSON_FILENAME}' through this endpoint")

        file_key = f"{prefix}/{clean_name}" if prefix else clean_name
        logger.info(f"🗑️ Deleting s3://{bucket_name}/{file_key}")
        _get_s3_client().delete_object(Bucket=bucket_name, Key=file_key)
        deleted.append(clean_name)
        logger.info(f"✅ Deleted s3://{bucket_name}/{file_key}")

    return deleted


def save_project_response_json(bucket_name: str, project_folder: str, response_data: dict) -> str:
    """
    Save the project overview response as response.json in the project folder.

    Returns:
        The S3 key where the file was saved
    """
    file_key = get_response_json_key(project_folder)
    body = json.dumps(response_data, indent=2, ensure_ascii=False).encode("utf-8")

    upload_file_to_s3(
        file_content=body,
        bucket_name=bucket_name,
        file_key=file_key,
        content_type="application/json",
    )
    logger.info(f"💾 Cached project overview saved to s3://{bucket_name}/{file_key}")
    return file_key


def get_project_response_json(bucket_name: str, project_folder: str) -> dict | None:
    """
    Load a cached project overview from response.json in S3.

    Returns:
        Parsed JSON dict if found, otherwise None
    """
    

    file_key = get_response_json_key(project_folder)
    try:
        logger.info(f"📂 Loading cached project overview from s3://{bucket_name}/{file_key}")
        raw = download_file_from_s3(bucket_name, file_key)
        logger.info(f"✅ Cached project overview loaded from s3://{bucket_name}/{file_key}")
        return json.loads(raw.decode("utf-8"))
    except ClientError as e:
        if e.response["Error"]["Code"] in ("NoSuchKey", "404"):
            logger.info(f"📭 No cached project overview found at s3://{bucket_name}/{file_key}")
            return None
        raise


def read_document_from_s3(bucket_name, file_key):
    try:
        logger.info(
            f"📥 Downloading the {file_key} file from the bucket {bucket_name}...")
        response = _get_s3_client().get_object(Bucket=bucket_name, Key=file_key)
        file_content = response['Body'].read()
        file_extension = file_key.lower().split('.')[-1]

        return _process_file_content(file_extension, file_content)

    except Exception as e:
        logger.error(
            f"❌ Error while reading {file_key} from bucket {bucket_name}: {str(e)}")
        raise


def download_file_from_s3(bucket_name: str, file_key: str) -> bytes:
    """
    Download raw file bytes from S3 without any content processing.
    Useful for passing documents directly to Amazon Textract or other services.

    Args:
        bucket_name: Name of the S3 bucket
        file_key: The key (path) of the file in S3

    Returns:
        Raw file content as bytes

    Raises:
        Exception: If the download fails
    """
    try:
        logger.info(f"📥 Downloading raw file {file_key} from bucket {bucket_name}...")
        response = _get_s3_client().get_object(Bucket=bucket_name, Key=file_key)
        file_content = response['Body'].read()
        logger.info(f"✅ Successfully downloaded {file_key} ({len(file_content)} bytes)")
        return file_content
    except Exception as e:
        logger.error(f"Error downloading {file_key} from bucket {bucket_name}: {str(e)}")
        raise


def upload_file_to_s3(file_content, bucket_name, file_key, content_type=None):
    """
    Upload a file to an S3 bucket.

    Args:
        file_content: The content of the file to upload (bytes or file-like object)
        bucket_name: Name of the S3 bucket
        file_key: The key (path) where the file will be stored in S3
        content_type: Optional MIME type of the file

    Returns:
        dict: The response from S3 upload operation

    Raises:
        Exception: If the upload fails
    """
    try:
        logger.info(f"📤 Uploading file to {bucket_name}/{file_key}...")

        # Prepare upload parameters
        upload_args = {
            'Bucket': bucket_name,
            'Key': file_key,
            'Body': file_content
        }

        # Add content type if provided
        if content_type:
            upload_args['ContentType'] = content_type

        # Upload the file
        response = _get_s3_client().put_object(**upload_args)

        logger.info(
            f"✅ File successfully uploaded to {bucket_name}/{file_key}")
        return response

    except Exception as e:
        logger.error(
            f"❌ Error uploading file to {bucket_name}/{file_key}: {str(e)}")
        raise
